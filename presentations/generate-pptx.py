#!/usr/bin/env python3
"""
Searchboost — SEO Säljpresentation PowerPoint Generator
Dark theme med Searchboost-branding (pink/cyan/dark)

Stödjer:
  - Enkel data (auto-analys)
  - Utökad data (Trello-analys med high_priority, medium_priority, long_term, conclusion)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import sys
import os

# Searchboost colors
PINK = RGBColor(0xE9, 0x1E, 0x8C)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
GREEN = RGBColor(0x00, 0xE6, 0x76)
PURPLE = RGBColor(0x7C, 0x4D, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
DARK_CARD = RGBColor(0x1E, 0x29, 0x3B)
DARK_CARD2 = RGBColor(0x16, 0x21, 0x33)
RED = RGBColor(0xEF, 0x44, 0x44)
YELLOW = RGBColor(0xF5, 0x9E, 0x0B)
ORANGE = RGBColor(0xF9, 0x73, 0x16)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=16, color=WHITE, line_spacing=Pt(8), font_name='Calibri'):
    """Add multiple lines of text with consistent formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line_data, tuple):
            text, clr, sz, bld = line_data
        else:
            text, clr, sz, bld = line_data, color, font_size, False
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.space_after = line_spacing
    return txBox


def add_pink_bar(slide, left, top, width, height=Inches(0.06)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PINK
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, color=DARK_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card_with_border(slide, left, top, width, height, fill_color=DARK_CARD, border_color=PINK, border_width=Pt(2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    return shape


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, icon="•", icon_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{icon} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
    return txBox


def add_numbered_list(slide, left, top, width, height, items, font_size=15, color=WHITE, num_color=PINK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i+1}. {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
    return txBox


def add_section_header(slide, text):
    """Standard section header (pink label at top)."""
    add_pink_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
                 text, font_size=14, color=PINK, bold=True)


def generate_presentation(data, output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    company = data.get('company_name', 'Företag')
    domain = data.get('domain', 'example.se')
    score = data.get('score', 0)
    summary = data.get('summary', '')
    conclusion = data.get('conclusion', '')
    issues = data.get('issues', [])
    cost = data.get('cost_estimate', {})
    page_speed = data.get('page_speed', {})

    # Utökad data (från Trello-analys)
    high_priority = data.get('high_priority', [])
    medium_priority = data.get('medium_priority', [])
    long_term = data.get('long_term', [])

    critical = [i for i in issues if i.get('severity') == 'high']
    warnings = [i for i in issues if i.get('severity') == 'medium']
    info_issues = [i for i in issues if i.get('severity') in ('low', 'info')]

    is_comprehensive = bool(high_priority or medium_priority or long_term)

    # PageSpeed data
    mobile_speed = None
    desktop_speed = None
    if page_speed:
        mobile = page_speed.get('mobile', {})
        desktop = page_speed.get('desktop', {})
        if mobile:
            mobile_speed = mobile.get('score', None)
        if desktop:
            desktop_speed = desktop.get('score', None)

    # ═══════════════════════════════════════
    # SLIDE 1: Framsida
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_pink_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))

    add_text_box(slide, Inches(1), Inches(1.2), Inches(6), Inches(0.6),
                 'SEARCHBOOST', font_size=22, color=PINK, bold=True)

    add_text_box(slide, Inches(1), Inches(2.3), Inches(11), Inches(1.2),
                 'SEO-Analys & Åtgärdsplan', font_size=48, color=WHITE, bold=True)

    add_text_box(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.8),
                 company, font_size=40, color=CYAN, bold=True)

    add_text_box(slide, Inches(1), Inches(4.4), Inches(11), Inches(0.5),
                 domain, font_size=22, color=LIGHT_GRAY)

    add_pink_bar(slide, Inches(1), Inches(5.2), Inches(3), Inches(0.06))

    add_text_box(slide, Inches(1), Inches(5.5), Inches(6), Inches(0.5),
                 'Mikael Larsson — Medierådgivare', font_size=16, color=LIGHT_GRAY)
    add_text_box(slide, Inches(1), Inches(5.9), Inches(6), Inches(0.5),
                 'mikael@searchboost.nu | 0760-19 49 05', font_size=14, color=LIGHT_GRAY)

    # Searchboost branding bottom right
    add_text_box(slide, Inches(9), Inches(6.6), Inches(4), Inches(0.5),
                 'searchboost.nu', font_size=14, color=PINK, bold=True, alignment=PP_ALIGN.RIGHT)

    # ═══════════════════════════════════════
    # SLIDE 2: Score + Sammanfattning + PageSpeed
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, 'Nulägesanalys')

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
                 f'Analys av {domain}', font_size=32, color=WHITE, bold=True)

    # Score circle
    score_color = GREEN if score >= 70 else YELLOW if score >= 40 else RED
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1.8), Inches(2.8), Inches(2.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = DARK_CARD
    circle.line.color.rgb = score_color
    circle.line.width = Pt(6)

    add_text_box(slide, Inches(1), Inches(2.4), Inches(2.8), Inches(1),
                 str(score), font_size=64, color=score_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(2.8), Inches(0.4),
                 '/100', font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.8), Inches(2.8), Inches(0.4),
                 'SEO-poäng', font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Stats boxes
    stats = [
        (str(len(critical)), 'Kritiska', RED),
        (str(len(warnings)), 'Varningar', YELLOW),
        (str(len(issues)), 'Totalt', CYAN),
    ]
    for idx, (val, label, clr) in enumerate(stats):
        x = Inches(4.5) + Inches(idx * 2.5)
        add_card(slide, x, Inches(1.8), Inches(2.2), Inches(1.4))
        add_text_box(slide, x, Inches(1.95), Inches(2.2), Inches(0.7),
                     val, font_size=42, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(2.6), Inches(2.2), Inches(0.4),
                     label, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # PageSpeed boxes (if available)
    if mobile_speed is not None or desktop_speed is not None:
        ps_y = Inches(3.5)
        if mobile_speed is not None:
            m_color = GREEN if mobile_speed >= 70 else YELLOW if mobile_speed >= 40 else RED
            add_card(slide, Inches(4.5), ps_y, Inches(2.2), Inches(1.4))
            add_text_box(slide, Inches(4.5), ps_y + Inches(0.15), Inches(2.2), Inches(0.7),
                         str(mobile_speed), font_size=38, color=m_color, bold=True, alignment=PP_ALIGN.CENTER)
            add_text_box(slide, Inches(4.5), ps_y + Inches(0.8), Inches(2.2), Inches(0.4),
                         'Mobil PageSpeed', font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

        if desktop_speed is not None:
            d_color = GREEN if desktop_speed >= 70 else YELLOW if desktop_speed >= 40 else RED
            add_card(slide, Inches(7), ps_y, Inches(2.2), Inches(1.4))
            add_text_box(slide, Inches(7), ps_y + Inches(0.15), Inches(2.2), Inches(0.7),
                         str(desktop_speed), font_size=38, color=d_color, bold=True, alignment=PP_ALIGN.CENTER)
            add_text_box(slide, Inches(7), ps_y + Inches(0.8), Inches(2.2), Inches(0.4),
                         'Desktop PageSpeed', font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Summary box
    add_card(slide, Inches(4.5), Inches(5.2), Inches(8), Inches(2))
    summary_text = summary[:500] if summary else 'Ingen sammanfattning tillgänglig.'
    add_text_box(slide, Inches(4.8), Inches(5.4), Inches(7.5), Inches(1.7),
                 summary_text, font_size=13, color=WHITE)

    # ═══════════════════════════════════════
    # SLIDE 3: Identifierade problem — Kritiska
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, 'Identifierade problem')

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
                 'Vad hindrar er synlighet i Google?', font_size=32, color=WHITE, bold=True)

    # Critical issues — left column
    if critical:
        add_card(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(5.2))
        add_text_box(slide, Inches(0.9), Inches(2.0), Inches(5.3), Inches(0.4),
                     'Kritiska problem', font_size=18, color=RED, bold=True)
        add_pink_bar(slide, Inches(0.9), Inches(2.5), Inches(2), Inches(0.04))
        crit_items = [i.get('description', '')[:100] for i in critical]
        add_bullet_list(slide, Inches(0.9), Inches(2.7), Inches(5.3), Inches(4),
                        crit_items, font_size=13, color=WHITE, icon="▸")

    # Warnings — right column
    if warnings:
        add_card(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(5.2))
        add_text_box(slide, Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.4),
                     'Varningar', font_size=18, color=YELLOW, bold=True)
        add_pink_bar(slide, Inches(7.1), Inches(2.5), Inches(2), Inches(0.04))
        warn_items = [i.get('description', '')[:100] for i in warnings]
        add_bullet_list(slide, Inches(7.1), Inches(2.7), Inches(5.3), Inches(4),
                        warn_items, font_size=13, color=WHITE, icon="▸")

    # ═══════════════════════════════════════
    # SLIDE 4: Åtgärdsplan — Hög prioritet
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, 'Åtgärdsplan')

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
                 'Hög prioritet — Måste åtgärdas', font_size=32, color=WHITE, bold=True)

    if high_priority:
        # Use detailed high_priority list from Trello
        add_card(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(5.2))
        add_text_box(slide, Inches(0.9), Inches(2.0), Inches(11), Inches(0.4),
                     'Dessa åtgärder ger störst effekt på er synlighet', font_size=16, color=CYAN)
        add_numbered_list(slide, Inches(0.9), Inches(2.6), Inches(11), Inches(4.2),
                          high_priority, font_size=17, color=WHITE)
    elif critical:
        # Fallback: use critical issue descriptions
        add_card(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(5.2))
        add_text_box(slide, Inches(0.9), Inches(2.0), Inches(11), Inches(0.4),
                     'Dessa problem kräver omedelbar åtgärd', font_size=16, color=CYAN)
        fixes = [i.get('fix', i.get('description', ''))[:120] for i in critical]
        add_numbered_list(slide, Inches(0.9), Inches(2.6), Inches(11), Inches(4.2),
                          fixes, font_size=17, color=WHITE)

    # ═══════════════════════════════════════
    # SLIDE 5: Medelprioritering & Långsiktig
    # ═══════════════════════════════════════
    if medium_priority or long_term:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_section_header(slide, 'Åtgärdsplan — Fortsättning')

        add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
                     'Medel & Långsiktig prioritet', font_size=32, color=WHITE, bold=True)

        if medium_priority:
            card_height = Inches(3) if long_term else Inches(5.2)
            add_card(slide, Inches(0.6), Inches(1.8), Inches(5.8), card_height)
            add_text_box(slide, Inches(0.9), Inches(2.0), Inches(5.3), Inches(0.4),
                         'Medel prioritet', font_size=18, color=YELLOW, bold=True)
            add_pink_bar(slide, Inches(0.9), Inches(2.5), Inches(2), Inches(0.04))
            add_bullet_list(slide, Inches(0.9), Inches(2.7), Inches(5.3), Inches(2.5 if long_term else 4),
                            medium_priority[:8], font_size=14, color=WHITE, icon="▸")

        if long_term:
            lt_y = Inches(1.8)
            lt_x = Inches(6.8) if medium_priority else Inches(0.6)
            lt_w = Inches(5.8) if medium_priority else Inches(12)
            card_height = Inches(3) if medium_priority else Inches(5.2)

            add_card(slide, lt_x, lt_y, lt_w, card_height)
            add_text_box(slide, lt_x + Inches(0.3), lt_y + Inches(0.2), lt_w - Inches(0.6), Inches(0.4),
                         'Långsiktigt', font_size=18, color=PURPLE, bold=True)
            add_pink_bar(slide, lt_x + Inches(0.3), lt_y + Inches(0.7), Inches(2), Inches(0.04))
            add_bullet_list(slide, lt_x + Inches(0.3), lt_y + Inches(0.9), lt_w - Inches(0.6), Inches(2),
                            long_term[:6], font_size=14, color=WHITE, icon="▸")

    elif warnings or info_issues:
        # Fallback for simple data
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_section_header(slide, 'Åtgärdsplan — Fortsättning')

        add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
                     'Medel & Långsiktig prioritet', font_size=32, color=WHITE, bold=True)

        add_card(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(5.2))
        add_text_box(slide, Inches(0.9), Inches(2.0), Inches(11), Inches(0.4),
                     'Förbättringar som stärker er SEO ytterligare', font_size=16, color=CYAN)
        fixes2 = [i.get('fix', i.get('description', ''))[:120] for i in (warnings + info_issues)[:8]]
        add_bullet_list(slide, Inches(0.9), Inches(2.6), Inches(11), Inches(4.2),
                        fixes2, font_size=15, color=WHITE, icon="▸")

    # ═══════════════════════════════════════
    # SLIDE 6: Slutsats / Sammanfattning (om utökad data)
    # ═══════════════════════════════════════
    if conclusion:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_section_header(slide, 'Sammanfattning')

        add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
                     'Vår bedömning', font_size=32, color=WHITE, bold=True)

        add_card(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(4.5))
        add_text_box(slide, Inches(1), Inches(2.1), Inches(11.2), Inches(3.8),
                     conclusion, font_size=17, color=WHITE)

        # Mini summary stats at bottom
        bottom_y = Inches(5.6)
        mini_stats = [
            (str(len(critical)), 'Kritiska', RED),
            (str(len(warnings)), 'Varningar', YELLOW),
            (str(len(high_priority)), 'Hög prio åtgärder', PINK),
            (str(len(medium_priority)), 'Medel prio', CYAN),
        ]
        for idx, (val, label, clr) in enumerate(mini_stats):
            x = Inches(0.8) + Inches(idx * 3.1)
            add_card(slide, x, bottom_y, Inches(2.8), Inches(1.1))
            add_text_box(slide, x, bottom_y + Inches(0.1), Inches(2.8), Inches(0.5),
                         val, font_size=30, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
            add_text_box(slide, x, bottom_y + Inches(0.6), Inches(2.8), Inches(0.3),
                         label, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════
    # SLIDE 7: Investering / Pris
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, 'Investering')

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
                 'Vad ni får', font_size=32, color=WHITE, bold=True)

    # What's included
    included = [
        'Komplett SEO-analys & teknisk genomgång',
        'Optimering av alla meta-titlar & beskrivningar',
        'Schema markup (strukturerad data)',
        'Indexeringsstrategi & canonical-hantering',
        'Intern länkning & breadcrumbs',
        'Automatisk veckooptimering av era sidor',
        'Veckorapporter med resultat & framsteg',
        'Personlig kontakt & löpande rådgivning',
    ]

    add_card(slide, Inches(0.6), Inches(1.7), Inches(7.2), Inches(5.3))
    add_bullet_list(slide, Inches(0.9), Inches(1.9), Inches(6.8), Inches(5),
                    included, font_size=16, color=WHITE, icon="✓", icon_color=GREEN)

    # Price box
    cost_amount = cost.get('amount', 5000) if isinstance(cost, dict) else cost
    cost_tier = cost.get('tier', 'standard') if isinstance(cost, dict) else 'standard'

    tier_labels = {
        'basic': 'BASIC',
        'standard': 'STANDARD',
        'premium': 'PREMIUM',
    }
    tier_label = tier_labels.get(cost_tier, cost_tier.upper())

    # Main price card
    add_card_with_border(slide, Inches(8.3), Inches(2.0), Inches(4.3), Inches(4), border_color=PINK)

    add_text_box(slide, Inches(8.3), Inches(2.2), Inches(4.3), Inches(0.4),
                 tier_label, font_size=16, color=PINK, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(8.3), Inches(2.8), Inches(4.3), Inches(1),
                 f'{cost_amount:,} kr'.replace(',', ' '), font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(8.3), Inches(3.8), Inches(4.3), Inches(0.4),
                 '/månad exkl. moms', font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_pink_bar(slide, Inches(9), Inches(4.4), Inches(2.6), Inches(0.04))

    # Cost breakdown if available
    breakdown = cost.get('breakdown', []) if isinstance(cost, dict) else []
    if breakdown:
        for bdi, bd_item in enumerate(breakdown[:3]):
            name = bd_item.get('name', '')
            price = bd_item.get('price', '')
            add_text_box(slide, Inches(8.6), Inches(4.7) + Inches(bdi * 0.35), Inches(2.5), Inches(0.3),
                         name, font_size=11, color=LIGHT_GRAY)
            add_text_box(slide, Inches(11), Inches(4.7) + Inches(bdi * 0.35), Inches(1.3), Inches(0.3),
                         price, font_size=11, color=WHITE, alignment=PP_ALIGN.RIGHT)

    # ═══════════════════════════════════════
    # SLIDE 8: CTA
    # ═══════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_pink_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                 'Redo att synas i Google?', font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(3), Inches(9), Inches(0.8),
                 f'Vi kan börja optimera {domain} redan denna vecka.',
                 font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # CTA button shape
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(4.2), Inches(4.3), Inches(0.9))
    btn.fill.solid()
    btn.fill.fore_color.rgb = PINK
    btn.line.fill.background()
    add_text_box(slide, Inches(4.5), Inches(4.3), Inches(4.3), Inches(0.7),
                 'Starta idag', font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_pink_bar(slide, Inches(4), Inches(5.8), Inches(5.3), Inches(0.04))

    add_text_box(slide, Inches(1), Inches(6.1), Inches(11), Inches(0.4),
                 'Mikael Larsson | mikael@searchboost.nu | 0760-19 49 05 | searchboost.nu',
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
                 'SEARCHBOOST', font_size=18, color=PINK, bold=True, alignment=PP_ALIGN.CENTER)

    # Save
    prs.save(output_path)
    print(f'Presentation sparad: {output_path}')
    return output_path


if __name__ == '__main__':
    # Read data from command line arg (JSON file) or stdin
    if len(sys.argv) > 1:
        with open(sys.argv[1]) as f:
            data = json.load(f)
    else:
        data = json.load(sys.stdin)

    output = data.get('output_path', 'presentation.pptx')

    # Resolve relative paths based on script location
    if not os.path.isabs(output):
        script_dir = os.path.dirname(os.path.abspath(__file__))
        output = os.path.join(script_dir, output)

    # Ensure output directory exists
    os.makedirs(os.path.dirname(output), exist_ok=True)

    generate_presentation(data, output)
