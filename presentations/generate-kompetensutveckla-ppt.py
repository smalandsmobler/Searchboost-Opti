#!/usr/bin/env python3
"""
Kompetensutveckla.se — Ny sajt-struktur presentation
Searchboost dark theme, ämnesbaserad struktur
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Searchboost-färger
PINK    = RGBColor(0xE9, 0x1E, 0x8C)
CYAN    = RGBColor(0x00, 0xD4, 0xFF)
GREEN   = RGBColor(0x00, 0xE6, 0x76)
PURPLE  = RGBColor(0x7C, 0x4D, 0xFF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0x94, 0xA3, 0xB8)
DARK    = RGBColor(0x0F, 0x17, 0x2A)
CARD    = RGBColor(0x1E, 0x29, 0x3B)
CARD2   = RGBColor(0x16, 0x21, 0x33)
RED     = RGBColor(0xEF, 0x44, 0x44)
YELLOW  = RGBColor(0xF5, 0x9E, 0x0B)
ORANGE  = RGBColor(0xF9, 0x73, 0x16)

W = Inches(13.333)
H = Inches(7.5)


def bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def bar(slide, l, t, w, h=Inches(0.07), color=PINK):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()


def card(slide, l, t, w, h, color=CARD, border=None, bw=Pt(2)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()


def rect(slide, l, t, w, h, color=CARD):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()


def txt(slide, l, t, w, h, text, sz=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = color; p.font.bold = bold
    p.font.name = 'Calibri'; p.alignment = align


def bullets(slide, l, t, w, h, items, sz=15, color=WHITE, icon='▸', spacing=Pt(8)):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f'{icon}  {item}'
        p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.name = 'Calibri'; p.space_after = spacing


def header(slide, label):
    bar(slide, 0, 0, W, Inches(0.08))
    txt(slide, Inches(0.7), Inches(0.35), Inches(10), Inches(0.45),
        label, sz=13, color=PINK, bold=True)
    txt(slide, Inches(9), Inches(0.35), Inches(4), Inches(0.45),
        'searchboost.nu', sz=12, color=LGRAY, align=PP_ALIGN.RIGHT)


# ───────────────────────────────────────────
def make():
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H
    blank = prs.slide_layouts[6]

    # ── SLIDE 1: Framsida ──────────────────
    s = prs.slides.add_slide(blank); bg(s)
    bar(s, 0, 0, W, Inches(0.09))

    txt(s, Inches(1), Inches(1.1), Inches(6), Inches(0.55),
        'SEARCHBOOST', sz=22, color=PINK, bold=True)

    txt(s, Inches(1), Inches(2.0), Inches(11), Inches(1.3),
        'Ny webbplatsstruktur', sz=52, color=WHITE, bold=True)
    txt(s, Inches(1), Inches(3.35), Inches(11), Inches(1.0),
        'Kompetensutveckla.se', sz=44, color=CYAN, bold=True)

    txt(s, Inches(1), Inches(4.45), Inches(11), Inches(0.55),
        'Ämnesbaserad struktur — fler sök, fler kursdeltagare', sz=20, color=LGRAY)

    bar(s, Inches(1), Inches(5.3), Inches(3.5), Inches(0.06))
    txt(s, Inches(1), Inches(5.55), Inches(7), Inches(0.45),
        'Mikael Larsson — Medierådgivare', sz=15, color=LGRAY)
    txt(s, Inches(1), Inches(5.95), Inches(7), Inches(0.45),
        'mikael@searchboost.nu  |  0760-19 49 05', sz=13, color=LGRAY)

    # ── SLIDE 2: Nuläge — 3 problem ───────
    s = prs.slides.add_slide(blank); bg(s)
    header(s, 'Nuläge')
    txt(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.65),
        'Tre saker som kostar er trafik idag', sz=30, color=WHITE, bold=True)

    problems = [
        ('971 trasiga\ninterna länkar', '971', RED,
         'URL-strukturen ändrades utan 301-redirects. Google tappar förtroendet för sajten.'),
        ('200+ gratisresurser\nosynliga för Google', '200+', YELLOW,
         'Checklistor, mallar, AFS-texter ligger gömda — ingen söktrafik trots värdefullt innehåll.'),
        ('12+ toppkategorier\nkaotisk struktur', '12+', ORANGE,
         'Samma ämne på 3 ställen. Google förstår inte vad sajten handlar om.'),
    ]

    for i, (title, val, clr, desc) in enumerate(problems):
        x = Inches(0.5 + i * 4.25)
        card(s, x, Inches(1.8), Inches(4.0), Inches(5.2), border=clr)
        # Stor siffra
        txt(s, x, Inches(2.0), Inches(4.0), Inches(1.4),
            val, sz=62, color=clr, bold=True, align=PP_ALIGN.CENTER)
        bar(s, x + Inches(0.3), Inches(3.4), Inches(3.4), Inches(0.05), color=clr)
        # Titel
        txt(s, x + Inches(0.25), Inches(3.6), Inches(3.5), Inches(0.9),
            title, sz=17, color=WHITE, bold=True)
        # Beskrivning
        txt(s, x + Inches(0.25), Inches(4.55), Inches(3.5), Inches(2.2),
            desc, sz=14, color=LGRAY)

    # ── SLIDE 3: Dagens vs ny struktur ────
    s = prs.slides.add_slide(blank); bg(s)
    header(s, 'Strukturjämförelse')
    txt(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.65),
        'Från kaos till tydlighet', sz=30, color=WHITE, bold=True)

    # Vänster — Dagens
    card(s, Inches(0.5), Inches(1.65), Inches(5.9), Inches(5.45), color=CARD2, border=RED)
    txt(s, Inches(0.5), Inches(1.75), Inches(5.9), Inches(0.5),
        'Idag — Leveransformat som kategori', sz=15, color=RED, bold=True, align=PP_ALIGN.CENTER)
    bar(s, Inches(0.8), Inches(2.25), Inches(5.3), Inches(0.04), color=RED)
    today = [
        'Webb-utbildningar  (kategori)',
        'Fysiska kurser  (kategori)',
        'Lärarledda kurser  (kategori)',
        '',
        'Ledarskap & chefsutbildning  (överlapp)',
        'Ledarskapsutbildning  (överlapp)',
        'Ledarskap  (överlapp)',
        '',
        '12+ toppkategorier, alla blandat',
        'Kurskatalog och kunskapsbank ihopblandat',
    ]
    bullets(s, Inches(0.8), Inches(2.4), Inches(5.4), Inches(4.4),
            today, sz=13, color=LGRAY, icon='✗', spacing=Pt(5))

    # Höger — Ny
    card(s, Inches(6.9), Inches(1.65), Inches(5.9), Inches(5.45), color=CARD2, border=GREEN)
    txt(s, Inches(6.9), Inches(1.75), Inches(5.9), Inches(0.5),
        'Ny — Ämne som kategori', sz=15, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    bar(s, Inches(7.2), Inches(2.25), Inches(5.3), Inches(0.04), color=GREEN)
    new_items = [
        '/utbildningar/  (EN toppkategori)',
        '  /ledarskap/',
        '  /hr-personalhantering/',
        '  /arbetsmiljo/',
        '  /ekonomi-juridik/',
        '  /digitalt-it/',
        '/kunskapsbank/  (NY toppkategori)',
        '  /checklistor/',
        '  /mallar/',
        'Format (webb/fysisk) = filter/tagg, ej kategori',
    ]
    bullets(s, Inches(7.2), Inches(2.4), Inches(5.4), Inches(4.4),
            new_items, sz=13, color=WHITE, icon='✓', spacing=Pt(5))

    # Pil mitt emellan
    txt(s, Inches(6.1), Inches(3.8), Inches(0.85), Inches(0.8),
        '→', sz=40, color=PINK, bold=True, align=PP_ALIGN.CENTER)

    # ── SLIDE 4: Den nya strukturen (visuell) ──
    s = prs.slides.add_slide(blank); bg(s)
    header(s, 'Ny struktur — ämnesbaserad')
    txt(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.65),
        'Så ser den nya sajten ut', sz=30, color=WHITE, bold=True)

    # Rotnod
    card(s, Inches(4.8), Inches(1.7), Inches(3.7), Inches(0.75), border=PINK)
    txt(s, Inches(4.8), Inches(1.82), Inches(3.7), Inches(0.5),
        'kompetensutveckla.se', sz=16, color=PINK, bold=True, align=PP_ALIGN.CENTER)

    # Nivå 2 — 2 toppkategorier
    toppkat = [
        (Inches(1.0), '/utbildningar/', CYAN),
        (Inches(8.5), '/kunskapsbank/', GREEN),
    ]
    for lx, label, clr in toppkat:
        card(s, lx, Inches(2.85), Inches(3.2), Inches(0.7), border=clr)
        txt(s, lx, Inches(2.97), Inches(3.2), Inches(0.5),
            label, sz=15, color=clr, bold=True, align=PP_ALIGN.CENTER)
        # Linje upp till rot
        bar(s, Inches(6.65), Inches(2.45), Inches(0.04), Inches(0.4), color=LGRAY)

    # Utbildnings-underkategorier
    utb_cats = [
        '/ledarskap/',
        '/hr-personal/',
        '/arbetsmiljo/',
        '/ekonomi-juridik/',
        '/digitalt-it/',
        '/halsa-friskvard/',
    ]
    for i, cat in enumerate(utb_cats):
        cx = Inches(0.35 + i * 2.1)
        card(s, cx, Inches(4.0), Inches(1.9), Inches(0.6), color=CARD2, border=CYAN)
        txt(s, cx, Inches(4.1), Inches(1.9), Inches(0.4),
            cat, sz=10, color=CYAN, align=PP_ALIGN.CENTER)

    bar(s, Inches(2.5), Inches(3.55), Inches(0.04), Inches(0.45), color=LGRAY)

    # Kunskapsbankens underkategorier
    kb_cats = ['/checklistor/', '/mallar/', '/AFS-texter/', '/guider/']
    for i, cat in enumerate(kb_cats):
        cx = Inches(7.8 + i * 1.35)
        card(s, cx, Inches(4.0), Inches(1.2), Inches(0.6), color=CARD2, border=GREEN)
        txt(s, cx, Inches(4.1), Inches(1.2), Inches(0.4),
            cat, sz=10, color=GREEN, align=PP_ALIGN.CENTER)

    bar(s, Inches(10.1), Inches(3.55), Inches(0.04), Inches(0.45), color=LGRAY)

    # Filter-taggar (format)
    txt(s, Inches(0.5), Inches(4.95), Inches(12), Inches(0.35),
        'Format-taggar (visas som filter på listningar):', sz=13, color=LGRAY)
    tags = ['Webb-utbildning', 'Fysisk kurs', 'Lärarledd', 'På plats', 'Halvdag', 'Heldag']
    for i, tag in enumerate(tags):
        card(s, Inches(0.5 + i * 2.1), Inches(5.35), Inches(1.95), Inches(0.5), color=PURPLE)
        txt(s, Inches(0.5 + i * 2.1), Inches(5.43), Inches(1.95), Inches(0.35),
            tag, sz=12, color=WHITE, align=PP_ALIGN.CENTER)

    # Not om kunskapshubbar
    txt(s, Inches(0.5), Inches(6.1), Inches(12.5), Inches(0.5),
        'Varje ämnes-kategori (/ledarskap/, /arbetsmiljo/ etc.) får en pillar page — rankad landningssida som samlar alla kurser och resurser inom ämnet.',
        sz=12, color=LGRAY)

    # ── SLIDE 5: Nyckelordsmöjligheter ────
    s = prs.slides.add_slide(blank); bg(s)
    header(s, 'Nyckelordsmöjligheter')
    txt(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.65),
        '8 300 månliga sökningar — ni har redan innehållet', sz=30, color=WHITE, bold=True)

    # Stor siffra
    card(s, Inches(0.5), Inches(1.7), Inches(3.2), Inches(2.6))
    txt(s, Inches(0.5), Inches(1.9), Inches(3.2), Inches(1.3),
        '8 300', sz=56, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(0.5), Inches(3.15), Inches(3.2), Inches(0.4),
        'sök/mån totalt', sz=15, color=LGRAY, align=PP_ALIGN.CENTER)
    txt(s, Inches(0.5), Inches(3.55), Inches(3.2), Inches(0.5),
        'ni har redan innehåll för', sz=13, color=LGRAY, align=PP_ALIGN.CENTER)

    # Nyckelords-kategorier
    kw_data = [
        ('Ledarskap & chefsutbildning', '3 400 sök/mån', PINK, [
            'chefsutbildning  1 600/mån',
            'ledarskapsutbildning  880/mån',
            'ledarskap utbildning  390/mån',
        ]),
        ('Arbetsmiljö & HR', '2 900 sök/mån', YELLOW, [
            'arbetsmiljöutbildning  590/mån',
            'hr utbildning  480/mån',
            'skyddsombud utbildning  390/mån',
        ]),
        ('Ekonomi & Juridik', '2 000 sök/mån', GREEN, [
            'löneadministratör utbildning  880/mån',
            'arbetsrättslig utbildning  390/mån',
            'bokföring utbildning  480/mån',
        ]),
    ]

    for i, (cat, vol, clr, kws) in enumerate(kw_data):
        x = Inches(4.0 + i * 3.1)
        card(s, x, Inches(1.7), Inches(2.95), Inches(5.4), border=clr)
        txt(s, x + Inches(0.2), Inches(1.9), Inches(2.6), Inches(0.5),
            cat, sz=14, color=clr, bold=True)
        txt(s, x + Inches(0.2), Inches(2.45), Inches(2.6), Inches(0.4),
            vol, sz=17, color=WHITE, bold=True)
        bar(s, x + Inches(0.2), Inches(2.9), Inches(2.4), Inches(0.04), color=clr)
        bullets(s, x + Inches(0.2), Inches(3.05), Inches(2.6), Inches(3.5),
                kws, sz=12, color=LGRAY, icon='▸', spacing=Pt(8))

    # Nuläge vs potential
    txt(s, Inches(0.5), Inches(6.25), Inches(12.5), Inches(0.5),
        'Ni har idag ~52 klick/mån. Med rätt struktur och pillar pages: mål 500–1 000 klick/mån inom 6–12 månader.',
        sz=13, color=LGRAY)

    # ── SLIDE 6: Vad ingår ────────────────
    s = prs.slides.add_slide(blank); bg(s)
    header(s, 'Vad ingår i omläggningen')
    txt(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.65),
        'Komplett omläggning — inget ni behöver göra själva', sz=30, color=WHITE, bold=True)

    left_items = [
        ('Kategoristruktur-omläggning', 'Komplett omstrukturering från leveransformat till ämnesbaserad', CYAN),
        ('URL-struktur optimering', 'Nya SEO-vänliga URL:er för alla 600+ sidor', PINK),
        ('Taxonomi & filter-taggar', 'Format (Webb/Fysisk/Lärarledd) som filter istället för kategori', PURPLE),
        ('Kunskapshub-uppsättning', 'Ny toppkategori med pillar pages för checklistor, mallar, AFS mm', GREEN),
    ]
    right_items = [
        ('Intern länkning', 'Alla sidor länkade i rätt hierarki, breadcrumbs korrekt', CYAN),
        ('EduAdmin-konfiguration', 'Synka kurskategorier med den nya strukturen', YELLOW),
        ('QA & testning', 'Allt testas i staging — ni godkänner innan det går live', ORANGE),
        ('Sökordsoptimerade landningssidor', 'Pillar pages per ämne — rankar på breda söktermer', GREEN),
    ]

    for i, (title, desc, clr) in enumerate(left_items):
        y = Inches(1.75 + i * 1.35)
        card(s, Inches(0.5), y, Inches(6.0), Inches(1.2), border=clr)
        txt(s, Inches(0.75), y + Inches(0.1), Inches(5.5), Inches(0.45),
            title, sz=15, color=clr, bold=True)
        txt(s, Inches(0.75), y + Inches(0.55), Inches(5.5), Inches(0.55),
            desc, sz=12, color=LGRAY)

    for i, (title, desc, clr) in enumerate(right_items):
        y = Inches(1.75 + i * 1.35)
        card(s, Inches(7.0), y, Inches(6.0), Inches(1.2), border=clr)
        txt(s, Inches(7.25), y + Inches(0.1), Inches(5.5), Inches(0.45),
            title, sz=15, color=clr, bold=True)
        txt(s, Inches(7.25), y + Inches(0.55), Inches(5.5), Inches(0.55),
            desc, sz=12, color=LGRAY)

    # ── SLIDE 7: ROI ──────────────────────
    s = prs.slides.add_slide(blank); bg(s)
    header(s, 'Avkastning på investering')
    txt(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.65),
        'Konservativ beräkning — vad ni kan förvänta er', sz=30, color=WHITE, bold=True)

    # 4 siffror
    roi_stats = [
        ('+2 000', 'extra besökare/mån', GREEN),
        ('20 st', 'extra bokningar/mån', CYAN),
        ('2 000 kr', 'snitt per bokning', YELLOW),
        ('480 000 kr', 'extra intäkt/år', PINK),
    ]
    for i, (val, label, clr) in enumerate(roi_stats):
        x = Inches(0.5 + i * 3.25)
        card(s, x, Inches(1.7), Inches(3.0), Inches(2.0), border=clr)
        txt(s, x, Inches(1.85), Inches(3.0), Inches(0.95),
            val, sz=38, color=clr, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x, Inches(2.82), Inches(3.0), Inches(0.45),
            label, sz=14, color=LGRAY, align=PP_ALIGN.CENTER)

    # Vad ni dessutom får
    card(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(3.0))
    txt(s, Inches(0.8), Inches(4.15), Inches(12), Inches(0.45),
        'Utöver intäkter', sz=16, color=CYAN, bold=True)
    bar(s, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.04), color=CYAN)

    bonus = [
        '971 trasiga interna länkar fixade — bättre användarupplevelse direkt',
        'Topical authority i Google — ni äger "arbetsmiljö utbildning"-nischen',
        'Skalbar struktur — nya utbildningar läggs enkelt till rätt ämne',
        'Kunskapshubbar som lead-generatorer — gratis innehåll driver betald kursbokning',
        'Bättre EduAdmin-integration — kurskatalog synkar med SEO-optimerad struktur',
    ]
    # Dela i 2 kolumner
    col1 = bonus[:3]; col2 = bonus[3:]
    bullets(s, Inches(0.8), Inches(4.75), Inches(6.0), Inches(2.0),
            col1, sz=13, color=WHITE, icon='✓', spacing=Pt(6))
    bullets(s, Inches(7.0), Inches(4.75), Inches(5.8), Inches(2.0),
            col2, sz=13, color=WHITE, icon='✓', spacing=Pt(6))

    # ── SLIDE 8: Tidplan ──────────────────
    s = prs.slides.add_slide(blank); bg(s)
    header(s, 'Tidplan')
    txt(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.65),
        'Klart på 3 veckor — allt i testmiljö först', sz=30, color=WHITE, bold=True)

    weeks = [
        ('Vecka 1', 'Analys & planering', CYAN, [
            'SEO-audit av hela sajten',
            'Komplett nyckelordsanalys',
            'Kartlägga alla kurser → nya kategorier',
            'Bygga redirect-lista (142+ URL:er)',
            'Teknisk genomgång',
        ]),
        ('Vecka 2', 'Bygg staging', PINK, [
            'Ny kategoristruktur i staging',
            'EduAdmin-konfiguration',
            'URL-flytt & pillar pages',
            'Kunskapshub-uppsättning',
            'Breadcrumbs & intern länkning',
        ]),
        ('Vecka 3', 'QA & lansering', GREEN, [
            'QA & testning',
            'Ni godkänner staging',
            'Stegvis lansering i produktion',
            'Monitorering & finjustering',
            'Ni äger resultatet',
        ]),
    ]

    for i, (week, title, clr, items) in enumerate(weeks):
        x = Inches(0.5 + i * 4.27)
        card(s, x, Inches(1.65), Inches(4.05), Inches(5.4), border=clr)
        # Vecka-nummer
        card(s, x + Inches(0.2), Inches(1.85), Inches(0.95), Inches(0.95), color=clr)
        txt(s, x + Inches(0.2), Inches(1.93), Inches(0.95), Inches(0.8),
            str(i+1), sz=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x + Inches(1.3), Inches(1.9), Inches(2.7), Inches(0.45),
            week, sz=13, color=clr, bold=True)
        txt(s, x + Inches(1.3), Inches(2.3), Inches(2.7), Inches(0.4),
            title, sz=14, color=WHITE, bold=True)
        bar(s, x + Inches(0.2), Inches(2.85), Inches(3.5), Inches(0.04), color=clr)
        bullets(s, x + Inches(0.2), Inches(3.0), Inches(3.6), Inches(3.7),
                items, sz=13, color=WHITE, icon='▸', spacing=Pt(7))

    txt(s, Inches(0.5), Inches(7.15), Inches(12.5), Inches(0.35),
        'Ni godkänner i staging innan något går live. Noll risk.',
        sz=13, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    # ── SLIDE 9: Pris ─────────────────────
    s = prs.slides.add_slide(blank); bg(s)
    header(s, 'Investering')
    txt(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.65),
        'Vad ingår och vad det kostar', sz=30, color=WHITE, bold=True)

    # Pristabell vänster
    card(s, Inches(0.5), Inches(1.7), Inches(7.0), Inches(4.8), color=CARD2)
    price_rows = [
        ('Kategoristruktur-omläggning', '7 500 kr'),
        ('URL-struktur optimering', '4 500 kr'),
        ('Taxonomi & tagg-optimering', '5 250 kr'),
        ('Breadcrumb-struktur', '3 750 kr'),
        ('Kunskapshub-uppsättning', 'ingår'),
        ('EduAdmin-konfiguration', 'ingår'),
        ('QA & staging-test', 'ingår'),
    ]
    txt(s, Inches(0.8), Inches(1.85), Inches(6.5), Inches(0.45),
        'Ingår i totalpaket', sz=14, color=CYAN, bold=True)
    bar(s, Inches(0.8), Inches(2.3), Inches(6.2), Inches(0.04), color=CYAN)

    for i, (name, price) in enumerate(price_rows):
        y = Inches(2.45 + i * 0.52)
        price_clr = GREEN if price == 'ingår' else LGRAY
        txt(s, Inches(0.8), y, Inches(5.0), Inches(0.45), name, sz=13, color=WHITE)
        txt(s, Inches(5.6), y, Inches(1.6), Inches(0.45), price, sz=13, color=price_clr,
            bold=(price == 'ingår'), align=PP_ALIGN.RIGHT)
        if i < len(price_rows) - 1:
            bar(s, Inches(0.8), y + Inches(0.48), Inches(6.2), Inches(0.02), color=CARD)

    # Priskort höger
    card(s, Inches(8.0), Inches(1.7), Inches(4.8), Inches(4.8), border=PINK, bw=Pt(3))
    txt(s, Inches(8.0), Inches(2.0), Inches(4.8), Inches(0.5),
        'TOTALPAKET', sz=16, color=PINK, bold=True, align=PP_ALIGN.CENTER)
    bar(s, Inches(8.5), Inches(2.5), Inches(3.8), Inches(0.06), color=PINK)
    txt(s, Inches(8.0), Inches(2.65), Inches(4.8), Inches(1.3),
        '23 250 kr', sz=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(8.0), Inches(3.95), Inches(4.8), Inches(0.45),
        'exkl. moms, engångspris', sz=14, color=LGRAY, align=PP_ALIGN.CENTER)
    bar(s, Inches(8.5), Inches(4.45), Inches(3.8), Inches(0.04), color=LGRAY)
    txt(s, Inches(8.0), Inches(4.6), Inches(4.8), Inches(0.45),
        'Allt inkluderat — inga dolda kostnader', sz=13, color=GREEN, align=PP_ALIGN.CENTER)
    txt(s, Inches(8.0), Inches(5.1), Inches(4.8), Inches(0.45),
        'Leverans 3 veckor från godkännande', sz=13, color=LGRAY, align=PP_ALIGN.CENTER)
    txt(s, Inches(8.0), Inches(5.6), Inches(4.8), Inches(0.45),
        'Ni godkänner innan det går live', sz=13, color=LGRAY, align=PP_ALIGN.CENTER)

    # Momstext
    txt(s, Inches(0.5), Inches(6.6), Inches(12.5), Inches(0.4),
        'Alla priser exkl. moms.  Totalpaket inkluderar kunskapshub och EduAdmin-konfiguration utan extra kostnad.',
        sz=11, color=LGRAY)

    # ── SLIDE 10: Nästa steg / CTA ────────
    s = prs.slides.add_slide(blank); bg(s)
    bar(s, 0, 0, W, Inches(0.09))

    txt(s, Inches(1), Inches(1.4), Inches(11), Inches(1.1),
        'Redo att ta täten\ni kursbranschen?', sz=46, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    txt(s, Inches(2), Inches(3.0), Inches(9), Inches(0.7),
        'Vi kan starta SEO-audit och nyckelordsanalys redan denna vecka — utan EduAdmin-access.',
        sz=18, color=LGRAY, align=PP_ALIGN.CENTER)

    # 3 steg
    steps = [
        ('1', 'Ni godkänner förslaget', PINK),
        ('2', 'Vi startar audit & planering', CYAN),
        ('3', 'Ni godkänner staging — lansering', GREEN),
    ]
    for i, (num, step_txt, clr) in enumerate(steps):
        x = Inches(1.5 + i * 3.7)
        card(s, x, Inches(4.0), Inches(3.4), Inches(0.9), border=clr)
        txt(s, x + Inches(0.15), Inches(4.08), Inches(0.6), Inches(0.7),
            num, sz=22, color=clr, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.8), Inches(4.1), Inches(2.5), Inches(0.65),
            step_txt, sz=14, color=WHITE)

    bar(s, Inches(4.5), Inches(5.35), Inches(4.3), Inches(0.06))
    txt(s, Inches(1), Inches(5.6), Inches(11), Inches(0.45),
        'Mikael Larsson  |  mikael@searchboost.nu  |  0760-19 49 05',
        sz=15, color=LGRAY, align=PP_ALIGN.CENTER)
    txt(s, Inches(1), Inches(6.05), Inches(11), Inches(0.45),
        'SEARCHBOOST', sz=20, color=PINK, bold=True, align=PP_ALIGN.CENTER)

    # Spara
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       'output', 'kompetensutveckla-ny-struktur.pptx')
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print(f'Sparad: {out}')
    return out


if __name__ == '__main__':
    make()
