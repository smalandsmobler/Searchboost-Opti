#!/usr/bin/env python3
"""
Searchboost Opti — Systemförklaring PPT
Genererar presentations/output/searchboost-systemforklaring-2026.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# ── Färger (Searchboost dark theme) ──
BG_DARK     = RGBColor(0x0A, 0x0A, 0x0F)
BG_CARD     = RGBColor(0x14, 0x14, 0x1E)
SB_PINK     = RGBColor(0xE9, 0x1E, 0x8C)
SB_CYAN     = RGBColor(0x00, 0xD4, 0xFF)
SB_GREEN    = RGBColor(0x00, 0xE6, 0x76)
SB_PURPLE   = RGBColor(0x7C, 0x4D, 0xFF)
SB_ORANGE   = RGBColor(0xFF, 0x6B, 0x35)
TEXT_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MUTED  = RGBColor(0x8A, 0x8A, 0xA8)
BORDER      = RGBColor(0x2A, 0x2A, 0x3E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]  # Blank


def set_bg(slide, color=BG_DARK):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(1)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, font_size=Pt(18), bold=False,
             color=TEXT_WHITE, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_arrow(slide, x1, y1, x2, y2, color=SB_CYAN):
    """Enkel linje med pilspets (connector)."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # STRAIGHT = 1
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    # Lägg till pilhuvud
    ln = connector.line._ln
    tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'none')
    headEnd = etree.SubElement(ln, qn('a:headEnd'))
    headEnd.set('type', 'arrow')
    headEnd.set('w', 'med')
    headEnd.set('len', 'med')
    return connector


# ═══════════════════════════════════════════════════
# SLIDE 1 — Framsida
# ═══════════════════════════════════════════════════
def slide_framsida(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    # Accent-linje vänster
    add_rect(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, fill_color=SB_PINK)

    # Logotyp-text
    add_text(slide, "SEARCHBOOST",
             Inches(0.4), Inches(0.35), Inches(6), Inches(0.6),
             font_size=Pt(13), bold=True, color=SB_PINK)

    # Huvud-rubrik
    add_text(slide, "Hur systemet fungerar",
             Inches(0.4), Inches(1.5), Inches(9), Inches(1.5),
             font_size=Pt(52), bold=True, color=TEXT_WHITE)

    # Underrubrik
    add_text(slide, "Searchboost Opti — Automatisk SEO-optimering från prospekt till rapport",
             Inches(0.4), Inches(3.1), Inches(10), Inches(0.7),
             font_size=Pt(20), color=TEXT_MUTED)

    # Datum
    add_text(slide, "Februari 2026",
             Inches(0.4), Inches(6.8), Inches(4), Inches(0.5),
             font_size=Pt(13), color=TEXT_MUTED)

    # Dekorativ box höger
    add_rect(slide, Inches(10.5), Inches(1.2), Inches(2.6), Inches(4.8),
             fill_color=BG_CARD, line_color=SB_CYAN, line_width=Pt(1))

    for i, (label, val, color) in enumerate([
        ("Kunder", "10 st", SB_CYAN),
        ("Opt/mån", "~450 st", SB_GREEN),
        ("Lambda", "4 st", SB_PURPLE),
        ("Uptime", "99.9 %", SB_ORANGE),
    ]):
        yy = Inches(1.5) + i * Inches(1.0)
        add_text(slide, label,
                 Inches(10.7), yy, Inches(2.2), Inches(0.35),
                 font_size=Pt(11), color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        add_text(slide, val,
                 Inches(10.7), yy + Inches(0.3), Inches(2.2), Inches(0.45),
                 font_size=Pt(22), bold=True, color=color, align=PP_ALIGN.CENTER)

    return slide


# ═══════════════════════════════════════════════════
# SLIDE 2 — Systemöversikt (flöde)
# ═══════════════════════════════════════════════════
def slide_systemovershikt(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_text(slide, "Systemöversikt", Inches(0.4), Inches(0.3), Inches(12), Inches(0.6),
             font_size=Pt(32), bold=True, color=TEXT_WHITE)
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(3), Inches(0.04), fill_color=SB_PINK)

    steps = [
        ("1", "Prospekt",       "Lägg till\nvia Dashboard",  SB_CYAN),
        ("2", "SEO-Audit",      "Crawla WP\nHitta problem",  SB_PURPLE),
        ("3", "Åtgärdsplan",    "AI genererar\n3-månadsplan",SB_PINK),
        ("4", "Optimering",     "Claude skriver\nmetadata",  SB_GREEN),
        ("5", "Rapport",        "Veckomail\n+ Kundportal",   SB_ORANGE),
    ]

    box_w  = Inches(2.0)
    box_h  = Inches(2.4)
    gap    = Inches(0.45)
    start_x = Inches(0.45)
    y      = Inches(2.0)

    for i, (num, title, desc, color) in enumerate(steps):
        x = start_x + i * (box_w + gap)

        # Kortbakgrund
        add_rect(slide, x, y, box_w, box_h, fill_color=BG_CARD, line_color=color, line_width=Pt(1.5))

        # Nummer
        add_text(slide, num,
                 x, y + Inches(0.15), box_w, Inches(0.5),
                 font_size=Pt(30), bold=True, color=color, align=PP_ALIGN.CENTER)

        # Titel
        add_text(slide, title,
                 x, y + Inches(0.65), box_w, Inches(0.55),
                 font_size=Pt(17), bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

        # Beskrivning
        add_text(slide, desc,
                 x + Inches(0.1), y + Inches(1.25), box_w - Inches(0.2), Inches(0.9),
                 font_size=Pt(12), color=TEXT_MUTED, align=PP_ALIGN.CENTER)

        # Pil till nästa
        if i < len(steps) - 1:
            ax = x + box_w + Inches(0.05)
            ay = y + box_h / 2
            add_arrow(slide, ax, ay, ax + gap - Inches(0.1), ay, color=SB_CYAN)

    # Undertext
    add_text(slide, "Hela kedjan kör autonomt — Mikael säljer, systemet levererar.",
             Inches(0.4), Inches(6.6), Inches(12), Inches(0.55),
             font_size=Pt(14), color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    return slide


# ═══════════════════════════════════════════════════
# SLIDE 3 — Datainsamling (Weekly Audit)
# ═══════════════════════════════════════════════════
def slide_datainsamling(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_text(slide, "Datainsamling", Inches(0.4), Inches(0.3), Inches(8), Inches(0.6),
             font_size=Pt(32), bold=True, color=TEXT_WHITE)
    add_text(slide, "Lambda: weekly-audit  —  Måndag 06:00 CET",
             Inches(0.4), Inches(0.95), Inches(8), Inches(0.45),
             font_size=Pt(15), color=SB_CYAN)
    add_rect(slide, Inches(0.4), Inches(1.38), Inches(3), Inches(0.04), fill_color=SB_CYAN)

    items = [
        (SB_CYAN,   "Crawlar WP REST API",    "Hämtar alla sidor + metadata per kund"),
        (SB_PURPLE, "Analyserar SEO-brister", "Title, H1, meta-desc, schema, canonical, alt-text"),
        (SB_GREEN,  "Prioriterar uppgifter",  "Impact-score 1–10 avgör köordning"),
        (SB_PINK,   "Fyller seo_work_queue",  "BigQuery-tabell med status pending/done/error"),
        (SB_ORANGE, "GSC-data samlas",        "Sökord, positioner, klick — partitionerad per dag"),
    ]

    for i, (color, title, desc) in enumerate(items):
        y = Inches(1.8) + i * Inches(0.9)
        add_rect(slide, Inches(0.4), y, Inches(0.06), Inches(0.65), fill_color=color)
        add_text(slide, title,
                 Inches(0.65), y, Inches(5.5), Inches(0.38),
                 font_size=Pt(16), bold=True, color=TEXT_WHITE)
        add_text(slide, desc,
                 Inches(0.65), y + Inches(0.35), Inches(7), Inches(0.35),
                 font_size=Pt(12), color=TEXT_MUTED)

    # Höger: BigQuery-box
    add_rect(slide, Inches(9.0), Inches(1.6), Inches(3.9), Inches(4.8),
             fill_color=BG_CARD, line_color=SB_PURPLE, line_width=Pt(1))
    add_text(slide, "BigQuery — 10 tabeller",
             Inches(9.1), Inches(1.75), Inches(3.7), Inches(0.5),
             font_size=Pt(14), bold=True, color=SB_PURPLE, align=PP_ALIGN.CENTER)

    tables = [
        "customer_pipeline",
        "customer_keywords",
        "seo_work_queue",
        "seo_optimization_log",
        "action_plans",
        "weekly_reports",
        "gsc_daily_metrics",
        "ads_daily_metrics",
        "social_daily_metrics",
        "data_collection_log",
    ]
    for i, t in enumerate(tables):
        add_text(slide, f"• {t}",
                 Inches(9.2), Inches(2.25) + i * Inches(0.39), Inches(3.5), Inches(0.4),
                 font_size=Pt(11), color=TEXT_MUTED)

    return slide


# ═══════════════════════════════════════════════════
# SLIDE 4 — Automatisk optimering
# ═══════════════════════════════════════════════════
def slide_optimering(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_text(slide, "Automatisk optimering", Inches(0.4), Inches(0.3), Inches(10), Inches(0.6),
             font_size=Pt(32), bold=True, color=TEXT_WHITE)
    add_text(slide, "Lambda: autonomous-optimizer  —  Var 6:e timme",
             Inches(0.4), Inches(0.95), Inches(10), Inches(0.45),
             font_size=Pt(15), color=SB_GREEN)
    add_rect(slide, Inches(0.4), Inches(1.38), Inches(3.5), Inches(0.04), fill_color=SB_GREEN)

    # Flöde: kö → Claude → WP → logg
    steps = [
        ("seo_work_queue",   "Hämtar nästa\nuppgift (prio)",    SB_PURPLE),
        ("Claude Haiku",     "Genererar\noptimerad text",       SB_PINK),
        ("WordPress API",    "Skriver title,\ndesc, schema, H1",SB_CYAN),
        ("BigQuery",         "Loggar utfört\narbete",           SB_GREEN),
        ("Trello DONE",      "Skapar kort\ni DONE-listan",      SB_ORANGE),
    ]

    box_w  = Inches(2.0)
    box_h  = Inches(2.0)
    gap    = Inches(0.45)
    start_x = Inches(0.45)
    y = Inches(2.0)

    for i, (title, desc, color) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        add_rect(slide, x, y, box_w, box_h, fill_color=BG_CARD, line_color=color, line_width=Pt(1.5))
        add_text(slide, title,
                 x, y + Inches(0.15), box_w, Inches(0.55),
                 font_size=Pt(14), bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, desc,
                 x + Inches(0.1), y + Inches(0.75), box_w - Inches(0.2), Inches(0.9),
                 font_size=Pt(12), color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            ax = x + box_w + Inches(0.05)
            ay = y + box_h / 2
            add_arrow(slide, ax, ay, ax + gap - Inches(0.1), ay, color=SB_GREEN)

    # Regler-box under
    add_rect(slide, Inches(0.4), Inches(4.35), Inches(12.5), Inches(2.65),
             fill_color=BG_CARD, line_color=BORDER, line_width=Pt(1))
    add_text(slide, "Säkerhetsregler",
             Inches(0.6), Inches(4.5), Inches(4), Inches(0.45),
             font_size=Pt(14), bold=True, color=SB_ORANGE)

    rules = [
        ("Budgettak",     "Basic: 15 opt/mån  |  Standard: 30  |  Premium: 50"),
        ("Inga nyckelord","Skippar kunden helt denna körning"),
        ("Ingen åtgärdsplan","Kör bara: title, H1, schema, thin content"),
        ("Redirect-sida", "Skippar sidan — ändrar aldrig redirect-destinationer"),
    ]
    for i, (label, desc) in enumerate(rules):
        col = 0 if i < 2 else 1
        row = i % 2
        x = Inches(0.6) + col * Inches(6.2)
        y2 = Inches(4.95) + row * Inches(0.7)
        add_text(slide, f"{label}:", x, y2, Inches(2.2), Inches(0.4),
                 font_size=Pt(12), bold=True, color=SB_ORANGE)
        add_text(slide, desc, x + Inches(2.1), y2, Inches(3.8), Inches(0.4),
                 font_size=Pt(12), color=TEXT_MUTED)

    return slide


# ═══════════════════════════════════════════════════
# SLIDE 5 — Veckolog
# ═══════════════════════════════════════════════════
def slide_veckolog(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_text(slide, "Veckolog & Viktor-dag", Inches(0.4), Inches(0.3), Inches(10), Inches(0.6),
             font_size=Pt(32), bold=True, color=TEXT_WHITE)
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(3.5), Inches(0.04), fill_color=SB_ORANGE)

    # Vänster: veckomail
    add_rect(slide, Inches(0.4), Inches(1.2), Inches(5.8), Inches(5.6),
             fill_color=BG_CARD, line_color=SB_CYAN, line_width=Pt(1))
    add_text(slide, "Veckomail — Måndag 08:00",
             Inches(0.6), Inches(1.35), Inches(5.4), Inches(0.5),
             font_size=Pt(16), bold=True, color=SB_CYAN)

    mail_items = [
        "Antal optimeringar denna vecka",
        "Per kund — vilka sidor förbättrades",
        "Nyckelordspositioner (GSC)",
        "Kommande veckas planerade arbete",
        "Sammanfattning i BigQuery: weekly_reports",
    ]
    for i, item in enumerate(mail_items):
        add_text(slide, f"→  {item}",
                 Inches(0.7), Inches(2.0) + i * Inches(0.75), Inches(5.2), Inches(0.55),
                 font_size=Pt(13), color=TEXT_MUTED)

    # Höger: Viktor-dag
    add_rect(slide, Inches(6.7), Inches(1.2), Inches(6.2), Inches(5.6),
             fill_color=BG_CARD, line_color=SB_ORANGE, line_width=Pt(1))
    add_text(slide, "Viktor-dag — Tisdag 09:00",
             Inches(6.9), Inches(1.35), Inches(5.8), Inches(0.5),
             font_size=Pt(16), bold=True, color=SB_ORANGE)

    viktor_items = [
        "Trello-kort skapas automatiskt varje tisdag",
        "Kunder med fel i arbets­kön",
        "Kunder utan nyckelord / åtgärdsplan",
        "Inaktiva kunder (0 opt denna vecka)",
        "Statistik: totalt antal opt per kund",
        "Viktor bockar av punkterna under dagen",
    ]
    for i, item in enumerate(viktor_items):
        add_text(slide, f"✓  {item}",
                 Inches(6.9), Inches(2.0) + i * Inches(0.7), Inches(5.8), Inches(0.55),
                 font_size=Pt(13), color=TEXT_MUTED)

    return slide


# ═══════════════════════════════════════════════════
# SLIDE 6 — Månadsrapport
# ═══════════════════════════════════════════════════
def slide_manadsrapport(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_text(slide, "Månadsrapport", Inches(0.4), Inches(0.3), Inches(10), Inches(0.6),
             font_size=Pt(32), bold=True, color=TEXT_WHITE)
    add_text(slide, "Vad kunden ser — och vad vi mäter",
             Inches(0.4), Inches(0.95), Inches(10), Inches(0.45),
             font_size=Pt(16), color=TEXT_MUTED)
    add_rect(slide, Inches(0.4), Inches(1.38), Inches(3), Inches(0.04), fill_color=SB_PINK)

    metrics = [
        ("Organisk trafik",    "+X klick/mån\nvs förra månaden",  SB_CYAN),
        ("Nyckelordpositioner","Genomsnittlig\nposition GSC",      SB_GREEN),
        ("Utförda opt",        "Antal sidor\nförbättrade",         SB_PURPLE),
        ("ROI-uppskattning",   "Trafikvärde\nvs månadskostnad",    SB_PINK),
        ("Nästa steg",         "Prioriterad plan\nkommande 30 dagar",SB_ORANGE),
        ("Budget-status",      "Använda opt\nav budget­tak",       SB_CYAN),
    ]

    for i, (title, desc, color) in enumerate(metrics):
        col = i % 3
        row = i // 3
        x = Inches(0.4) + col * Inches(4.3)
        y = Inches(1.7) + row * Inches(2.5)
        add_rect(slide, x, y, Inches(4.0), Inches(2.2),
                 fill_color=BG_CARD, line_color=color, line_width=Pt(1))
        add_rect(slide, x, y, Inches(4.0), Inches(0.06), fill_color=color)
        add_text(slide, title,
                 x + Inches(0.15), y + Inches(0.2), Inches(3.7), Inches(0.5),
                 font_size=Pt(15), bold=True, color=TEXT_WHITE)
        add_text(slide, desc,
                 x + Inches(0.15), y + Inches(0.75), Inches(3.7), Inches(1.1),
                 font_size=Pt(13), color=TEXT_MUTED)

    return slide


# ═══════════════════════════════════════════════════
# SLIDE 7 — Kundzon (portal)
# ═══════════════════════════════════════════════════
def slide_kundzon(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_text(slide, "Kundzon", Inches(0.4), Inches(0.3), Inches(10), Inches(0.6),
             font_size=Pt(32), bold=True, color=TEXT_WHITE)
    add_text(slide, "opti.searchboost.nu — vad kunden ser",
             Inches(0.4), Inches(0.95), Inches(10), Inches(0.45),
             font_size=Pt(16), color=SB_CYAN)
    add_rect(slide, Inches(0.4), Inches(1.38), Inches(2), Inches(0.04), fill_color=SB_CYAN)

    sections = [
        ("Prestanda",      "Klick, impressioner,\ngenomsnittlig pos.",  SB_CYAN),
        ("Sökord",         "Alla rankade sökord\nmed trend-pilar",       SB_GREEN),
        ("Optimeringar",   "Tidslinje: vad vi\ngjort och när",           SB_PURPLE),
        ("Åtgärdsplan",    "Nästa 3 månaders\nprioriterade uppgifter",   SB_PINK),
        ("AI-chatt",       "Fråga om sin data:\n'Varför minskade klick?'",SB_ORANGE),
    ]

    for i, (title, desc, color) in enumerate(sections):
        x = Inches(0.4)
        y = Inches(1.8) + i * Inches(1.0)
        add_rect(slide, x, y, Inches(0.06), Inches(0.75), fill_color=color)
        add_text(slide, title,
                 Inches(0.65), y, Inches(3.0), Inches(0.4),
                 font_size=Pt(16), bold=True, color=TEXT_WHITE)
        add_text(slide, desc,
                 Inches(0.65), y + Inches(0.38), Inches(5.5), Inches(0.45),
                 font_size=Pt(12), color=TEXT_MUTED)

    # Mock-skärm höger
    add_rect(slide, Inches(7.2), Inches(1.3), Inches(5.7), Inches(5.7),
             fill_color=RGBColor(0x10, 0x10, 0x18), line_color=SB_CYAN, line_width=Pt(1.5))

    # Faux header
    add_rect(slide, Inches(7.2), Inches(1.3), Inches(5.7), Inches(0.45), fill_color=BG_CARD)
    add_text(slide, "opti.searchboost.nu/portal",
             Inches(7.35), Inches(1.35), Inches(5.2), Inches(0.35),
             font_size=Pt(10), color=TEXT_MUTED)

    rows = [
        (SB_CYAN,   "Organisk trafik",   "+23 % vs förra månaden"),
        (SB_GREEN,  "Genomsnittlig pos.", "14.2 → 11.8"),
        (SB_PURPLE, "Opt. denna månad",  "34 st klara"),
        (SB_ORANGE, "Budget",            "34/50 använda"),
    ]
    for i, (color, label, val) in enumerate(rows):
        y2 = Inches(2.0) + i * Inches(0.85)
        add_rect(slide, Inches(7.35), y2, Inches(5.3), Inches(0.72),
                 fill_color=BG_CARD, line_color=color, line_width=Pt(0.8))
        add_text(slide, label,
                 Inches(7.5), y2 + Inches(0.1), Inches(3), Inches(0.3),
                 font_size=Pt(10), color=TEXT_MUTED)
        add_text(slide, val,
                 Inches(7.5), y2 + Inches(0.35), Inches(4.8), Inches(0.28),
                 font_size=Pt(13), bold=True, color=color)

    return slide


# ═══════════════════════════════════════════════════
# SLIDE 8 — Tekniskt flöde
# ═══════════════════════════════════════════════════
def slide_tekniskt_flode(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_text(slide, "Tekniskt flöde", Inches(0.4), Inches(0.3), Inches(10), Inches(0.6),
             font_size=Pt(32), bold=True, color=TEXT_WHITE)
    add_text(slide, "Alla system som pratar med varandra",
             Inches(0.4), Inches(0.95), Inches(10), Inches(0.45),
             font_size=Pt(16), color=TEXT_MUTED)
    add_rect(slide, Inches(0.4), Inches(1.38), Inches(3), Inches(0.04), fill_color=SB_PURPLE)

    nodes = [
        # (label,        x,           y,          color,     size)
        ("EC2\n(Express)", Inches(5.8), Inches(3.0), SB_PINK,   Inches(1.6), Inches(0.9)),
        ("BigQuery",       Inches(9.8), Inches(2.2), SB_PURPLE, Inches(1.5), Inches(0.7)),
        ("WordPress",      Inches(9.8), Inches(4.3), SB_CYAN,   Inches(1.5), Inches(0.7)),
        ("GSC",            Inches(5.8), Inches(5.5), SB_GREEN,  Inches(1.2), Inches(0.7)),
        ("Trello",         Inches(1.5), Inches(5.5), SB_ORANGE, Inches(1.2), Inches(0.7)),
        ("Kund\n(Portal)", Inches(1.5), Inches(2.2), SB_CYAN,   Inches(1.2), Inches(0.8)),
        ("Lambda\n×4",     Inches(3.4), Inches(1.3), SB_PURPLE, Inches(1.2), Inches(0.8)),
        ("Claude\nHaiku",  Inches(8.2), Inches(1.3), SB_PINK,   Inches(1.2), Inches(0.8)),
    ]

    centers = {}
    for (label, x, y, color, w, h) in nodes:
        add_rect(slide, x, y, w, h, fill_color=BG_CARD, line_color=color, line_width=Pt(2))
        add_text(slide, label, x, y + Inches(0.08), w, h - Inches(0.08),
                 font_size=Pt(12), bold=True, color=color, align=PP_ALIGN.CENTER)
        cx = x + w / 2
        cy = y + h / 2
        centers[label.split('\n')[0]] = (cx, cy)

    # Pilar EC2 ↔ partners
    connections = [
        ("EC2",      "BigQuery",  SB_PURPLE),
        ("EC2",      "WordPress", SB_CYAN),
        ("EC2",      "GSC",       SB_GREEN),
        ("EC2",      "Trello",    SB_ORANGE),
        ("EC2",      "Kund",      SB_CYAN),
        ("Lambda",   "EC2",       SB_PURPLE),
        ("Claude",   "EC2",       SB_PINK),
    ]
    for (a, b, color) in connections:
        if a in centers and b in centers:
            x1, y1 = centers[a]
            x2, y2 = centers[b]
            add_arrow(slide, x1, y1, x2, y2, color=color)

    return slide


# ═══════════════════════════════════════════════════
# SLIDE 9 — Kundstatus
# ═══════════════════════════════════════════════════
def slide_kundstatus(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    add_text(slide, "Kundstatus — februari 2026", Inches(0.4), Inches(0.3), Inches(12), Inches(0.6),
             font_size=Pt(32), bold=True, color=TEXT_WHITE)
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(4), Inches(0.04), fill_color=SB_GREEN)

    headers = ["Kund", "WP", "GSC", "Nyckelord", "Åtgärdsplan", "Status"]
    col_w = [Inches(2.4), Inches(0.9), Inches(0.9), Inches(1.4), Inches(1.8), Inches(1.8)]
    col_x = [Inches(0.3)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    # Header-rad
    add_rect(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(0.45), fill_color=BG_CARD)
    for j, (h, x) in enumerate(zip(headers, col_x)):
        add_text(slide, h, x + Inches(0.05), Inches(1.25), col_w[j] - Inches(0.05), Inches(0.38),
                 font_size=Pt(12), bold=True, color=TEXT_MUTED)

    customers = [
        ("searchboost",             "OK", "OK", "Ja", "Ja",   "Aktiv",   SB_GREEN),
        ("kompetensutveckla",       "OK", "OK", "Ja", "Ja",   "Aktiv",   SB_GREEN),
        ("ilmonte",                 "OK", "—",  "Ja", "Ja",   "Aktiv",   SB_CYAN),
        ("mobelrondellen",          "OK", "OK", "Ja", "Ja",   "Aktiv",   SB_GREEN),
        ("phvast",                  "OK", "OK", "Ja", "Ja",   "Aktiv",   SB_GREEN),
        ("tobler",                  "OK", "—",  "Ja", "Delvis","Aktiv",  SB_CYAN),
        ("traficator",              "OK", "—",  "Ja", "Delvis","Aktiv",  SB_CYAN),
        ("smalandskontorsmobler",   "—",  "OK", "Ja", "Nej",  "Delvis",  SB_ORANGE),
        ("wedosigns",               "—",  "—",  "Nej","Nej",  "Saknas",  SB_ORANGE),
        ("ferox",                   "—",  "—",  "Nej","Nej",  "Ej aktiv",SB_PINK),
    ]

    for i, (name, wp, gsc, kw, ap, status, color) in enumerate(customers):
        y = Inches(1.7) + i * Inches(0.48)
        if i % 2 == 0:
            add_rect(slide, Inches(0.3), y, Inches(12.7), Inches(0.45),
                     fill_color=RGBColor(0x16, 0x16, 0x22))
        row = [name, wp, gsc, kw, ap, status]
        for j, (val, x) in enumerate(zip(row, col_x)):
            c = color if j == len(row) - 1 else (SB_GREEN if val == "OK" or val == "Ja" else (SB_ORANGE if val in ("—","Delvis","Saknas") else TEXT_MUTED))
            add_text(slide, val, x + Inches(0.05), y + Inches(0.05), col_w[j] - Inches(0.05), Inches(0.38),
                     font_size=Pt(11), color=c)

    add_text(slide, "7 av 10 kunder kör automatisk optimering idag",
             Inches(0.3), Inches(6.9), Inches(8), Inches(0.45),
             font_size=Pt(13), color=TEXT_MUTED)

    return slide


# ═══════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════
def main():
    out_dir = os.path.join(os.path.dirname(__file__), "output")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "searchboost-systemforklaring-2026.pptx")

    prs = new_prs()

    print("Bygger slides...")
    slide_framsida(prs)
    print("  1/9 Framsida")
    slide_systemovershikt(prs)
    print("  2/9 Systemöversikt")
    slide_datainsamling(prs)
    print("  3/9 Datainsamling")
    slide_optimering(prs)
    print("  4/9 Automatisk optimering")
    slide_veckolog(prs)
    print("  5/9 Veckolog & Viktor-dag")
    slide_manadsrapport(prs)
    print("  6/9 Månadsrapport")
    slide_kundzon(prs)
    print("  7/9 Kundzon")
    slide_tekniskt_flode(prs)
    print("  8/9 Tekniskt flöde")
    slide_kundstatus(prs)
    print("  9/9 Kundstatus")

    prs.save(out_path)
    print(f"\nSparad: {out_path}")


if __name__ == "__main__":
    main()
